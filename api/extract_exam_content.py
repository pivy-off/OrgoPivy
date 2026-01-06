"""
Script to extract text content from exam study guide files
"""
import sys
from pathlib import Path
import json

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("Warning: python-docx not installed. Install with: pip install python-docx")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")

BASE_DIR = Path(__file__).resolve().parent.parent
EXAM_DIR = BASE_DIR / "exam"
OUTPUT_DIR = BASE_DIR / "api" / "storage" / "practice_exams"

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

def extract_docx_text(file_path: Path) -> str:
    """Extract text from .docx file"""
    if not HAS_DOCX:
        return f"[Error: python-docx not installed. Cannot read {file_path.name}]"
    
    try:
        doc = Document(file_path)
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return "\n".join(text_parts)
    except Exception as e:
        return f"[Error reading {file_path.name}: {str(e)}]"

def extract_pptx_text(file_path: Path) -> str:
    """Extract text from .pptx file"""
    if not HAS_PPTX:
        return f"[Error: python-pptx not installed. Cannot read {file_path.name}]"
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {slide_num} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        return f"[Error reading {file_path.name}: {str(e)}]"

def extract_doc_text(file_path: Path) -> str:
    """Extract text from .doc file (older Word format)"""
    # .doc files require different library (python-docx2txt or antiword)
    # For now, return a message
    return f"[Note: .doc files require additional tools. Please convert {file_path.name} to .docx or .txt]"

def process_exam_files():
    """Process all exam files and extract content"""
    if not EXAM_DIR.exists():
        print(f"Exam directory not found: {EXAM_DIR}")
        return
    
    results = {}
    
    for file_path in EXAM_DIR.iterdir():
        if file_path.is_file() and file_path.name != "desktop.ini":
            print(f"Processing: {file_path.name}")
            
            if file_path.suffix.lower() == ".docx":
                content = extract_docx_text(file_path)
            elif file_path.suffix.lower() == ".pptx":
                content = extract_pptx_text(file_path)
            elif file_path.suffix.lower() == ".doc":
                content = extract_doc_text(file_path)
            else:
                content = f"[Unsupported file type: {file_path.suffix}]"
            
            # Save as text file
            output_file = OUTPUT_DIR / f"{file_path.stem}.txt"
            output_file.write_text(content, encoding="utf-8")
            
            results[file_path.name] = {
                "extracted": len(content) > 0 and not content.startswith("["),
                "size": len(content),
                "output": str(output_file)
            }
            
            print(f"  -> Extracted {len(content)} characters to {output_file.name}")
    
    # Save summary
    summary_file = OUTPUT_DIR / "extraction_summary.json"
    summary_file.write_text(json.dumps(results, indent=2), encoding="utf-8")
    
    print(f"\nExtraction complete! Files saved to: {OUTPUT_DIR}")
    print(f"Summary saved to: {summary_file}")

if __name__ == "__main__":
    process_exam_files()
